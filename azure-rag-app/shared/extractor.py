# shared/extractor.py
import hashlib
import io
import json
import logging
import os
from typing import Dict, Any
import html as html_lib  # for escaping

# Azure Document Intelligence (Form Recognizer)
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential

# Fallback libraries
try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except Exception:
    Image = None
    pytesseract = None

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """
    Uses Azure Document Intelligence for supported formats (pdf, docx, pptx, images).
    Falls back to simple parsers for txt/json and to python-docx/python-pptx/pytesseract where available.
    """

    SUPPORTED_BY_DOC_INTELLIGENCE = {"pdf", "docx", "pptx", "png", "jpg", "jpeg", "tif", "tiff"}

    def __init__(self, form_recognizer_endpoint: str, form_recognizer_key: str, default_locale: str = "en"):
        self.form_recognizer_endpoint = form_recognizer_endpoint
        self.form_recognizer_key = form_recognizer_key
        self.default_locale = default_locale

        if form_recognizer_endpoint and form_recognizer_key:
            self.doc_client = DocumentAnalysisClient(
                endpoint=self.form_recognizer_endpoint,
                credential=AzureKeyCredential(self.form_recognizer_key),
            )
        else:
            self.doc_client = None
            logger.warning("Form recognizer endpoint/key not provided. Document Intelligence will be skipped.")

    @staticmethod
    def _stable_doc_id(blob_path: str, file_bytes: bytes) -> str:
        """
        Deterministic id for the document. Use blob path + length hash to detect changes.
        """
        h = hashlib.sha256()
        # combine path and length; you can extend with last_modified if available
        h.update(f"{blob_path}|{len(file_bytes)}".encode("utf-8"))
        return h.hexdigest()

    @staticmethod
    def _infer_topic_from_blob_path(blob_path: str) -> str:
        """
        Infer topic from blob path. Assumes structure like 'topic_name/filename.ext'.
        """
        try:
            parts = blob_path.split("/")
            if len(parts) >= 2:
                return parts[0]
            # fallback: parent directory or default
            return os.path.basename(os.path.dirname(blob_path)) or "unknown"
        except Exception:
            return "unknown"

    @staticmethod
    def _text_from_docx_bytes(file_bytes: bytes) -> str:
        if docx is None:
            raise RuntimeError("python-docx is not installed.")
        bio = io.BytesIO(file_bytes)
        document = docx.Document(bio)
        paragraphs = [p.text for p in document.paragraphs if p.text]
        return "\n".join(paragraphs)

    @staticmethod
    def _text_from_pptx_bytes(file_bytes: bytes) -> str:
        if Presentation is None:
            raise RuntimeError("python-pptx is not installed.")
        bio = io.BytesIO(file_bytes)
        prs = Presentation(bio)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

    @staticmethod
    def _text_from_txt_bytes(file_bytes: bytes) -> str:
        # assume UTF-8 with fallback to latin-1
        try:
            return file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            return file_bytes.decode("latin-1", errors="ignore")

    @staticmethod
    def _text_from_json_bytes(file_bytes: bytes) -> str:
        try:
            obj = json.loads(file_bytes.decode("utf-8"))
            # pretty-print JSON as text to include keys/values in index
            return json.dumps(obj, indent=2, ensure_ascii=False)
        except Exception:
            return file_bytes.decode("utf-8", errors="ignore")

    @staticmethod
    def _text_from_image_bytes_with_tesseract(file_bytes: bytes) -> str:
        if Image is None or pytesseract is None:
            raise RuntimeError("Pillow or pytesseract not installed.")
        bio = io.BytesIO(file_bytes)
        img = Image.open(bio)
        return pytesseract.image_to_string(img)
    
    
    @staticmethod
    def _text_from_pdf_bytes(file_bytes: bytes, image_ocr: bool = True) -> str:
        """Extract text from PDF bytes.

        Strategy (best-effort, progressive):
        1. Try PyMuPDF (fitz) fast text extraction.
        2. Fallback to pdfminer.six extract_text.
        3. If still empty and image_ocr=True, render pages to images and OCR via pytesseract (uses PyMuPDF or pdf2image).
        4. Final fallback: attempt UTF-8/L1 decode.
        """
        # 1) PyMuPDF (fast, preserves layout)
        try:
            import fitz  # PyMuPDF
            try:
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                texts = []
                for page in doc:
                    try:
                        texts.append(page.get_text("text") or "")
                    except Exception:
                        try:
                            texts.append(page.get_text() or "")
                        except Exception:
                            texts.append("")
                text = "\n".join(t for t in texts if t)
                if text.strip():
                    return text
            except Exception:
                logger.debug("PyMuPDF opened PDF but extraction failed.")
        except Exception:
            logger.debug("PyMuPDF not available for PDF extraction.")

        # 2) pdfminer.six
        try:
            from pdfminer.high_level import extract_text
            bio = io.BytesIO(file_bytes)
            try:
                text = extract_text(bio)
                if text and text.strip():
                    return text
            except Exception:
                logger.debug("pdfminer.extract_text failed to parse PDF.")
        except Exception:
            logger.debug("pdfminer.six not available for PDF extraction.")

        # 3) OCR fallback via pytesseract
        if image_ocr and Image is not None and pytesseract is not None:
            # Prefer rendering with PyMuPDF if available
            try:
                import fitz
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                ocr_texts = []
                for page in doc:
                    try:
                        pix = page.get_pixmap(dpi=200)
                        img = Image.open(io.BytesIO(pix.tobytes("png")))
                        ocr_texts.append(pytesseract.image_to_string(img))
                    except Exception as e:
                        logger.debug("PDF page render+OCR failed: %s", e)
                text = "\n".join(t for t in ocr_texts if t)
                if text.strip():
                    return text
            except Exception:
                logger.debug("PyMuPDF not available for rendering pages for OCR.")

            # Try pdf2image as a fallback if available
            try:
                from pdf2image import convert_from_bytes
                images = convert_from_bytes(file_bytes, dpi=200)
                ocr_texts = [pytesseract.image_to_string(img) for img in images]
                text = "\n".join(t for t in ocr_texts if t)
                if text.strip():
                    return text
            except Exception:
                logger.debug("pdf2image not available or failed for OCR fallback.")

        # 4) Last resort: try to decode as text
        try:
            return file_bytes.decode("utf-8")
        except Exception:
            return file_bytes.decode("latin-1", errors="ignore")
    

    def _extract_with_doc_intelligence(self, file_bytes: bytes) -> dict:
        """
        Use DocumentAnalysisClient to extract textual content, key-values, and tables (HTML).
        Returns a dict with:
        - content: plain text + appended HTML tables (optional)
        - metadata: kv pairs + DI metadata
        - tables: structured grid data for each table
        - tables_html: HTML strings for each table
        """
        if not self.doc_client:
            raise RuntimeError("DocumentAnalysisClient not configured.")

        poller = self.doc_client.begin_analyze_document("prebuilt-document", document=file_bytes)
        result = poller.result()

        # --- Plain text content ---
        content = getattr(result, "content", None)
        if not content:
            content_parts = []
            try:
                for page in getattr(result, "pages", []) or []:
                    if hasattr(page, "content") and page.content:
                        content_parts.append(page.content)
            except Exception:
                pass
            content = "\n".join(content_parts)

        # --- Key-value pairs ---
        metadata = {}
        try:
            kvs = getattr(result, "key_value_pairs", None)
            if kvs:
                for kv in kvs:
                    if getattr(kv, "key", None) and getattr(kv, "value", None):
                        key_text = kv.key.content.strip() if hasattr(kv.key, "content") else str(kv.key)
                        value_text = kv.value.content.strip() if hasattr(kv.value, "content") else str(kv.value)
                        safe_key = key_text.lower().replace(" ", "_")[:64]
                        metadata[safe_key] = value_text
        except Exception:
            pass

        # --- Tables: structured + HTML ---
        tables = []
        tables_html = []

        for t_idx, table in enumerate(getattr(result, "tables", []) or []):
            rows = getattr(table, "row_count", 0) or 0
            cols = getattr(table, "column_count", 0) or 0

            # Build 2D grid
            grid = [["" for _ in range(cols)] for _ in range(rows)]

            # Cells
            for cell in getattr(table, "cells", []) or []:
                r = getattr(cell, "row_index", None)
                c = getattr(cell, "column_index", None)
                text = getattr(cell, "content", "") or ""

                # Optional: handle spans if available in your SDK
                row_span = getattr(cell, "row_span", 1) or 1
                col_span = getattr(cell, "column_span", 1) or 1

                if r is not None and c is not None and 0 <= r < rows and 0 <= c < cols:
                    grid[r][c] = text
                    # NOTE: True HTML row/colspan rendering requires emitting <td rowspan=.. colspan=..>
                    # To keep this simple, we only set the top-left cell text and leave other spanned cells empty.
                    # If you need exact spanning, see the advanced renderer below.

            # Structured store
            tables.append({
                "index": t_idx,
                "row_count": rows,
                "column_count": cols,
                "cells": grid
            })

            # Determine header rows
            header_rows = getattr(table, "header_row_count", 0) or 0
            header_count = header_rows if header_rows > 0 else (1 if rows > 0 else 0)

            # HTML rendering (simple, without explicit row/colspan attributes)
            html_table = self._table_to_html(
                grid=grid,
                header_row_count=header_count,
                column_count=cols,
                table_class="di-table"
            )
            tables_html.append(html_table)

        # --- Merge HTML tables into content (optional) ---
        # If your consumer expects HTML, you can append these below. If it expects plain text, consider
        # returning tables_html separately and NOT appending.
        base_text = content or ""
        if tables_html:
            # separate with a simple horizontal rule
            html_joined = "\n<hr/>\n".join(tables_html)
            # If your content is plain text, you may want to keep them separate instead of appending.
            base_text = (base_text + "\n\n" + html_joined).strip()

        # --- DI metadata ---
        di_meta = {
            "di_model_id": "prebuilt-document",
            "di_pages": len(getattr(result, "pages", []) or []),
            "di_tables_found": len(getattr(result, "tables", []) or []),
        }
        metadata.update(di_meta)
        metadata["di_tables"] = tables
        metadata["di_tables_html"] = tables_html

        return {
            "content": base_text,
            "metadata": metadata,
            "tables": tables,
            "tables_html": tables_html,
        }


    def _table_to_html(self,grid, header_row_count=1, column_count=None, table_class="di-table"):
        """
        Convert a 2D grid (list[list[str]]) to an HTML table string.
        - header_row_count: number of rows at the top that are headers.
        - column_count: optional, helps render the header separator correctly if grid rows have ragged lengths.
        - table_class: CSS class for styling.
        """
        rows = len(grid)
        cols = column_count if column_count is not None else (max((len(r) for r in grid), default=0))

        def td(content):
            return f"<td>{html_lib.escape(content or '')}</td>"

        def th(content):
            return f"<th>{html_lib.escape(content or '')}</th>"

        parts = []
        parts.append(f'<table class="{html_lib.escape(table_class)}">')

        # THEAD
        if rows > 0 and cols > 0 and header_row_count > 0:
            parts.append("<thead>")
            for r in range(min(header_row_count, rows)):
                parts.append("<tr>")
                # header cells
                for c in range(cols):
                    cell_text = grid[r][c] if c < len(grid[r]) else ""
                    parts.append(th(cell_text))
                parts.append("</tr>")
            parts.append("</thead>")

        # TBODY
        parts.append("<tbody>")
        body_start = min(header_row_count, rows)
        for r in range(body_start, rows):
            parts.append("<tr>")
            for c in range(cols):
                cell_text = grid[r][c] if c < len(grid[r]) else ""
                parts.append(td(cell_text))
            parts.append("</tr>")
        parts.append("</tbody>")

        parts.append("</table>")
        return "".join(parts)


    def extract(self, file_bytes: bytes, file_type: str, blob_path: str = None) -> Dict[str, Any]:
        """
        Main entry point.

        :param file_bytes: raw bytes from the blob
        :param file_type: lower-case extension without dot (e.g., 'pdf', 'docx', 'txt', 'png')
        :param blob_path: original blob path (for topic inference and stable id)
        :returns: normalized document dict:
          {
            "id": "<document_id>",
            "topic": "<topic>",
            "source_path": "<blob_path>",
            "content": "<extracted text>",
            "metadata": { "file_type": "...", "title": "...", "created_on": "...", "author": "...", ... }
          }
        """
        file_type = (file_type or "").lower().strip()
        blob_path = blob_path or "unknown"
        doc_id = self._stable_doc_id(blob_path, file_bytes)
        topic = self._infer_topic_from_blob_path(blob_path)

        normalized = {
            "id": doc_id,
            "topic": topic,
            "source_path": blob_path,
            "content": "",
            "metadata": {
                "file_type": file_type,
                "title": os.path.basename(blob_path) if blob_path else None,
                "created_on": None,
                "author": None,
            },
        }

        # 1) Try Document Intelligence for supported formats
        if file_type in self.SUPPORTED_BY_DOC_INTELLIGENCE and self.doc_client:
            try:
                di = self._extract_with_doc_intelligence(file_bytes)
                normalized["content"] = di.get("content", "") or ""
                # merge di metadata into normalized.metadata
                for k, v in di.get("metadata", {}).items():
                    # avoid collisions with top-level fields
                    normalized["metadata"][f"di_{k}"] = v
                # return early if meaningful content
                if normalized["content"].strip():
                    logger.info(f"Document Intelligence extraction successful for {normalized}")
                    return normalized
                # else fallthrough to fallback
            except Exception as ex:
                logger.warning("Document Intelligence extraction failed for %s: %s", blob_path, ex)
                # continue to fallback extraction

        # 2) Fallback extraction depending on file type
        try:
            if file_type == "txt":
                normalized["content"] = self._text_from_txt_bytes(file_bytes)
            elif file_type == "json":
                normalized["content"] = self._text_from_json_bytes(file_bytes)
            elif file_type == "docx":
                if docx is not None:
                    normalized["content"] = self._text_from_docx_bytes(file_bytes)
                else:
                    # if python-docx not present, attempt naive decode as fallback
                    normalized["content"] = self._text_from_txt_bytes(file_bytes)
            elif file_type == "pptx":
                if Presentation is not None:
                    normalized["content"] = self._text_from_pptx_bytes(file_bytes)
                else:
                    normalized["content"] = self._text_from_txt_bytes(file_bytes)
            elif file_type in {"png", "jpg", "jpeg", "tif", "tiff"}:
                if Image is not None and pytesseract is not None:
                    try:
                        normalized["content"] = self._text_from_image_bytes_with_tesseract(file_bytes)
                    except Exception as e:
                        logger.warning("pytesseract failed: %s", e)
                        normalized["content"] = ""
                else:
                    normalized["content"] = ""
            elif file_type == "pdf":
                # Try to extract PDF text with common fallbacks (PyMuPDF, pdfminer, OCR)
                try:
                    normalized["content"] = self._text_from_pdf_bytes(file_bytes, image_ocr=True)
                except Exception as e:
                    logger.warning("PDF fallback extraction failed: %s", e)
                    normalized["content"] = ""
            else:
                # Unknown format - attempt best-effort decode
                normalized["content"] = self._text_from_txt_bytes(file_bytes)
        except Exception as ex:
            logger.exception("Fallback extraction failed for %s: %s", blob_path, ex)
            normalized["content"] = ""

        return normalized


